#!/usr/bin/env python3
"""Build the Pantomath / HORUS Hackatom 2026 deck from the template.

Strategy: open the 4-slide template, duplicate slides 2 and 3 to reach 10
slides, then rewrite text in place and embed sim screenshots.

Keeps original fonts, colors, and layout fixtures from the template.
"""
import copy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

ROOT = Path(__file__).parent
TEMPLATE = ROOT / "Pantomath Presentation Hackatom 2026.pptx"
OUTPUT = ROOT / "Pantomath HORUS Hackatom 2026.pptx"

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
RID_ATTRS = (f"{{{R_NS}}}id", f"{{{R_NS}}}embed", f"{{{R_NS}}}link")


def duplicate_slide(prs, src_idx):
    """Duplicate slide at src_idx; append to end; return the new slide."""
    src = prs.slides[src_idx]
    new_slide = prs.slides.add_slide(src.slide_layout)
    # strip default placeholders that came from the layout
    for shp in list(new_slide.shapes):
        sp = shp._element
        sp.getparent().remove(sp)
    # mirror rels (images, charts) and build rId map
    rid_map = {}
    for rel in src.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        if rel.is_external:
            new_rId = new_slide.part.relate_to(
                rel.target_ref, rel.reltype, is_external=True
            )
        else:
            new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rel.rId] = new_rId
    # deep-copy each shape, remap rIds
    for shp in src.shapes:
        el = copy.deepcopy(shp._element)
        for node in el.iter():
            for attr in RID_ATTRS:
                if attr in node.attrib and node.attrib[attr] in rid_map:
                    node.attrib[attr] = rid_map[node.attrib[attr]]
        new_slide.shapes._spTree.insert_element_before(el, "p:extLst")
    return new_slide


def move_slide(prs, old_idx, new_idx):
    sld_id_lst = prs.slides._sldIdLst
    slides = list(sld_id_lst)
    sld_id_lst.remove(slides[old_idx])
    sld_id_lst.insert(new_idx, slides[old_idx])


def shape_by_id(slide, shape_id):
    for shp in slide.shapes:
        if shp.shape_id == shape_id:
            return shp
    return None


def rewrite_table(table):
    """Replace the template table contents with our metrics."""
    rows = [
        ["metric",              "isolated",   "hivemind"],
        ["classification acc",  "82.7%",      "97.8%"],
        ["worker dose (µSv)",   "7131",       "4841"],
        ["drums completed",     "92",         "95"],
        ["ILW → LLW misroutes", "15",         "0"],
    ]
    # template table is 4 rows × 5 cols; we have 5 rows × 3 cols. We write
    # into the first 3 columns of each row that exists, and leave any extra
    # cells blank.
    n_rows = min(len(table.rows), len(rows))
    n_cols = min(len(table.columns), len(rows[0]))
    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            set_text(cell.text_frame, str(rows[r][c]))
        # blank any extra cells in this row
        for c in range(n_cols, len(table.columns)):
            cell = table.cell(r, c)
            set_text(cell.text_frame, "")
    # blank any extra rows
    for r in range(n_rows, len(table.rows)):
        for c in range(len(table.columns)):
            cell = table.cell(r, c)
            set_text(cell.text_frame, "")


def set_text(text_frame, new_text):
    """Replace all text in the frame with new_text. Multi-paragraph via '\n'.

    Preserves the *exact* rPr (run properties) of the original first run and
    the pPr (paragraph properties) of the original first paragraph, by
    deepcopy-ing the XML elements. This keeps theme colors, custom fonts,
    bullet styles, indentation, etc. — anything python-pptx's high-level
    Font API silently drops.
    """
    body = text_frame._txBody
    paragraphs = body.findall(qn("a:p"))
    src_rPr = None
    src_pPr = None
    if paragraphs:
        p0 = paragraphs[0]
        pPr = p0.find(qn("a:pPr"))
        if pPr is not None:
            src_pPr = copy.deepcopy(pPr)
        for r in p0.findall(qn("a:r")):
            rPr = r.find(qn("a:rPr"))
            if rPr is not None:
                src_rPr = copy.deepcopy(rPr)
                break
    # clear existing paragraphs
    for p in paragraphs:
        body.remove(p)
    # write one paragraph per line
    for line in new_text.split("\n"):
        p = etree.SubElement(body, qn("a:p"))
        if src_pPr is not None:
            p.append(copy.deepcopy(src_pPr))
        r = etree.SubElement(p, qn("a:r"))
        if src_rPr is not None:
            r.append(copy.deepcopy(src_rPr))
        t = etree.SubElement(r, qn("a:t"))
        t.text = line


# ---------- content ----------

# Slide 1 — Title
# Shape 47 is the big BOTTOM-LEFT slot (12" wide). Shape 48 is small under it
# (we widen below). Shape 49 is a small label ABOVE 47 (we widen below).
S1_BIG = "HORUS"
S1_SUBTITLE = "Hivemind for Onboard Radiological Understanding & Sorting"
S1_LABEL = "Egypt  ·  Team Pantomath"
S1_META = "Hackatom 2026"  # added in a new top-right text box

# Slide 2 — Problem (slide-2 layout)
S2_NUM = "01"
S2_SUB = "The problem"
S2_TITLE = "One mislabeled drum is a regulatory event."
S2_BODY = (
    "A nuclear facility produces tens of thousands of waste drums per year. "
    "Most are easy — VLLW gloves, LLW filters. The expensive ones are the "
    "edge cases: a low-activity-looking container that hides an actinide "
    "spike. Mislabel it as LLW, and you've put intermediate-level waste on "
    "the wrong shelf. The cost is regulatory, dosimetric, and human.*\n\n"
    "Today this work runs on slow, dose-limited human inspectors and "
    "isolated robotic scanners that don't talk to each other. Every new "
    "edge case has to be re-learned, drum by drum, person by person. The "
    "fleet has no memory."
)
S2_FOOT = "* IAEA GSG-1 classes: VLLW / LLW / ILW / HLW. An ILW drum sitting in an LLW bay is a dose event for any worker on inspection rounds."

# Slide 3 — Vision (slide-3 layout, two text columns)
S3_NUM = "02"
S3_SUB = "The vision"
S3_TITLE = "A living, thinking hivemind for radioactive waste."
S3_LEFT = (
    "We turn waste management and decontamination into a coordinated "
    "organism. Robotic arms, mobile scanners, fixed lab instruments — all "
    "of them sense, all of them act, all of them feed one central memory.\n\n"
    "Field instruments train against the lab. The lab pushes corrections "
    "back to the field in seconds. The first agent to see a new failure "
    "mode immunizes the entire fleet by the next shift.\n\n"
    "No drum waits on a human to re-derive a rule that the fleet already "
    "knows."
)
S3_RIGHT = (
    "AGGREGATE  —  every scan, every dose reading, every camera frame goes "
    "to one coordinator's ledger.\n\n"
    "LEARN  —  the coordinator derives one discriminating quantity — the "
    "actinide-spike threshold — from authoritative HPGe lab labels.\n\n"
    "PUSH  —  the threshold goes back to every agent. Every agent now "
    "classifies like the lab.\n\n"
    "DECIDE IN SECONDS  —  a tricky drum that would have been misrouted "
    "yesterday is caught at first contact today."
)

# Slide 4 — What we built (slide-2 layout)
S4_NUM = "03"
S4_SUB = "Proof of concept"
S4_TITLE = "A working facility-scale simulation. Runnable on a laptop."
S4_BODY = (
    "HORUS is not a slide deck. It is a multi-agent simulation of a "
    "simplified PUREX reprocessing line — Fuel Receipt → Shearing → "
    "Dissolution → Solvent Extraction → HLW Concentration → "
    "Solidification — with a fleet of five robotic agents sorting drums "
    "into VLLW / LLW / ILW / HLW storage.\n\n"
    "We model real gamma physics (Co-60, Cs-137, Am-241, U-235/238, "
    "Pu-239), real detector response (NaI 7% FWHM, HPGe 0.2% FWHM at 1332 "
    "keV), inverse-square dose accumulation, battery drain, and robot "
    "decontamination cycles.\n\n"
    "Two modes — ISOLATED (agents act alone) and HIVEMIND (agents share a "
    "coordinator) — run side by side on the same seed, on the same drums, "
    "to make the architectural claim falsifiable."
)
S4_FOOT = "* See the live demo running alongside this deck."

# Slide 5 — The Coordinator (slide-2 layout)
S5_NUM = "04"
S5_SUB = "The coordinator"
S5_TITLE = "One central memory. One number that matters."
S5_BODY = (
    "Every field report — gamma spectrum, camera frame, confidence, "
    "location, dose accumulated — lands in a single ledger. The "
    "coordinator's job is not deep learning. Its job is to derive one "
    "discriminating scalar — the actinide-spike threshold — from a handful "
    "of authoritative HPGe lab labels.\n\n"
    "When a new agent comes online, it inherits the threshold instantly. "
    "When the threshold drifts — a new feed composition, a new packaging "
    "vendor — the coordinator retrains and pushes a new snapshot in "
    "seconds.\n\n"
    "One number. One ledger. One source of truth. Auditable enough for a "
    "regulator to read line by line."
)
S5_FOOT = "* In the live demo: model version v0 → v2, actinide threshold derives from None → 0.045 over the first 20 minutes of operation."

# Slide 6 — The Agents (slide-3 layout, two columns)
S6_NUM = "05"
S6_SUB = "The field agents"
S6_TITLE = "NaI + camera + dose meter + the inherited threshold."
S6_LEFT = (
    "Each mobile agent carries:\n\n"
    "  •  A 2\"×2\" NaI(Tl) gamma spectrometer (~7% FWHM at 662 keV)\n"
    "  •  A lightweight camera + ML head for container packaging\n"
    "  •  An integrating dose meter for self-protection\n"
    "  •  A local copy of the coordinator's learned threshold\n\n"
    "Low-confidence calls and disagreements between sensors get flagged. "
    "Flagged drums travel to the fixed HPGe lab for authoritative "
    "re-classification — and those new labels feed the next coordinator "
    "retrain."
)
S6_RIGHT = (
    "Real robots in a hot cell don't run forever. We model that:\n\n"
    "  •  Battery drains during work, recharges in the bay\n"
    "  •  Cumulative dose accumulates inverse-square near hot items\n"
    "  •  At 5 mSv the agent must visit decon\n"
    "  •  At 50 mSv it permanently fails\n\n"
    "When an agent goes down, the hivemind's shared task queue routes its "
    "drums to surviving agents. Single-agent classification cannot recover "
    "from this — its work simply stops."
)

# Slide 7 — Split-second decision (slide-2 layout)
S7_NUM = "06"
S7_SUB = "The decisive moment"
S7_TITLE = "A masquerading ILW drum, caught in milliseconds."
S7_BODY = (
    "Watch the live demo. A drum arrives at the Cleanup-ops generation "
    "point. Its surface dose rate reads LLW. Its NaI gamma spectrum, by "
    "gross activity, says LLW. Computer vision says the packaging is LLW. "
    "Three sensors, one wrong answer.\n\n"
    "ISOLATED MODE  —  the agent obeys the rules engine. The drum is "
    "dropped at LLW storage. Workers on inspection rounds get "
    "inverse-square cooked. The mistake compounds.\n\n"
    "HIVEMIND MODE  —  the same agent runs the coordinator's learned "
    "threshold on the actinide photopeak windows. The threshold fires. "
    "The drum is escalated to ILW and routed via the HPGe lab for "
    "confirmation. The hivemind learned this exact trick yesterday, from "
    "one HPGe re-class, and every agent now knows it forever."
)
S7_FOOT = "* Press T in the live demo to inject a tricky drum and watch it happen in real time."

# Slide 8 — Results (slide-4 layout, table + charts)
S8_NUM = "07"
S8_SUB = "Results"
S8_TITLE = "82.7% → 97.8% accuracy.  15 → 0 ILW misroutes.  Dose down 32%.  Throughput flat."

# Slide 9 — Today vs production (slide-3 layout)
S9_NUM = "08"
S9_SUB = "Grounded, not promised"
S9_TITLE = "What runs today. What a pilot adds."
S9_LEFT = (
    "REAL TODAY (in the running simulation)\n\n"
    "  •  Real gamma physics: Co-60, Cs-137, Am-241, U-235/238, Pu-239 "
    "lines and intensities from public nuclear data\n"
    "  •  NaI / HPGe resolution, Poisson counting statistics, "
    "inverse-square dose and detector geometry\n"
    "  •  Two-tier classifier: auditable rules engine + one learned "
    "scalar\n"
    "  •  Federated retrain loop and per-agent snapshot delivery\n"
    "  •  Robot battery, integrated dose, decon, permanent-failure model\n"
    "  •  Side-by-side ISOLATED vs HIVEMIND runs to falsify the claim"
)
S9_RIGHT = (
    "PRODUCTION ROADMAP (what a pilot adds)\n\n"
    "  •  ROS 2 / DDS for actual robot comms (vs in-process Python)\n"
    "  •  OPC UA bridge to plant I&C systems\n"
    "  •  Data diode between safety-class and conventional sides "
    "(IEC 60709 / IEC 62859)\n"
    "  •  Redundant hot-standby coordinator\n"
    "  •  Safety-class qualification for any decision used in regulatory "
    "documentation (IEC 61513, IEEE 7-4.3.2, RG 1.152)\n"
    "  •  Full HPGe cryostat handling and ANSI N42.14 QA"
)

# Slide 10 — Thanks + open questions (slide-2 layout)
S10_NUM = "09"
S10_SUB = "Thanks"
S10_TITLE = "What HORUS opens up — and who we owe."
S10_BODY = (
    "Thanks to the Hackatom organizers for putting this problem in front "
    "of a hackathon, to the judges for the time, and to ENSDF / IAEA whose "
    "gamma tables and class boundaries we leaned on hard.\n\n"
    "HORUS shipped on one weekend with one trained discriminator. Open "
    "questions we'd love to argue about:\n\n"
    "  •  ADVERSARIAL DRUMS  —  what breaks when a packaging vendor "
    "changes geometry?\n"
    "  •  THRESHOLD DRIFT  —  how fast must we retrain across feed "
    "chemistries?\n"
    "  •  SINGLE-QUEEN FAILURE  —  smallest hot-standby protocol that "
    "survives a corrupted snapshot?\n"
    "  •  AUTHORITY BOUNDARY  —  where does advisory ML end and "
    "safety-class decision begin?\n\n"
    "Come find us."
)
S10_FOOT = "Team Pantomath  ·  Egypt  ·  ahmedfathy075@gmail.com  ·  try it:  python run_demo.py"


# ---------- build ----------

def main():
    prs = Presentation(str(TEMPLATE))
    assert len(prs.slides) == 4, f"expected 4 slides, got {len(prs.slides)}"

    # Snapshot the source slide indices BEFORE we add any duplicates.
    # original order: [title=0, slide2=1, slide3=2, slide4=3]
    SRC_TITLE = 0
    SRC_S2 = 1
    SRC_S3 = 2
    SRC_S4 = 3

    # Duplicate plan:
    #   slide-2 layout: need 5 total (problem, what we built, coordinator,
    #     split-second, ask). We have 1 original → need 4 dups.
    #   slide-3 layout: need 3 total (vision, agents, today/prod). We have 1
    #     original → need 2 dups.
    s2_dups = [duplicate_slide(prs, SRC_S2) for _ in range(4)]
    s3_dups = [duplicate_slide(prs, SRC_S3) for _ in range(2)]

    # After duplication, the slide list is:
    #   [0] title
    #   [1] slide2_orig
    #   [2] slide3_orig
    #   [3] slide4_orig
    #   [4..7] slide2 dups
    #   [8..9] slide3 dups
    # We want final order:
    #   [0] title
    #   [1] slide2_orig  → Problem
    #   [2] slide3_orig  → Vision
    #   [3] slide2 dup0  → What we built
    #   [4] slide2 dup1  → Coordinator
    #   [5] slide3 dup0  → Agents
    #   [6] slide2 dup2  → Split-second
    #   [7] slide4_orig  → Results
    #   [8] slide3 dup1  → Today vs production
    #   [9] slide2 dup3  → Ask
    # Reorder via sldIdLst manipulation.
    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    # detach
    for s in sld_ids:
        sld_id_lst.remove(s)
    # reattach in desired order
    desired = [
        sld_ids[0],  # title
        sld_ids[1],  # slide2_orig
        sld_ids[2],  # slide3_orig
        sld_ids[4],  # s2_dup0
        sld_ids[5],  # s2_dup1
        sld_ids[8],  # s3_dup0
        sld_ids[6],  # s2_dup2
        sld_ids[3],  # slide4_orig
        sld_ids[9],  # s3_dup1
        sld_ids[7],  # s2_dup3
    ]
    for s in desired:
        sld_id_lst.append(s)

    # ---------- fill content ----------

    # Slide 1 — title. Widen the two narrow boxes so content fits on one line,
    # then add a small top-right meta line.
    slide = prs.slides[0]
    big = shape_by_id(slide, 47)
    sub = shape_by_id(slide, 48)
    lbl = shape_by_id(slide, 49)
    set_text(big.text_frame, S1_BIG)
    set_text(sub.text_frame, S1_SUBTITLE)
    set_text(lbl.text_frame, S1_LABEL)
    # widen subtitle from 2.88" → 12" so it doesn't wrap into a tower
    sub.width = Inches(12)
    sub.height = Inches(0.6)
    # widen the country/label line so "Egypt · Team Pantomath" fits
    lbl.width = Inches(5)
    # add a small meta box at top-right with Hackatom 2026
    meta = slide.shapes.add_textbox(Inches(17.5), Inches(0.5), Inches(4), Inches(0.4))
    mp = meta.text_frame.paragraphs[0]
    mr = mp.add_run()
    mr.text = S1_META
    mr.font.name = "Ubuntu Mono"
    mr.font.size = Pt(14)
    mr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    mp.alignment = 2  # PP_ALIGN.RIGHT

    # Helper to fill a slide-2-layout-style slide.
    # Source shape ids (slide 2):
    #   79 = title placeholder ("long title for slide with text")
    #   80 = body placeholder (lorem ipsum)
    #   81 = page number ("02")
    #   82 = subtitle ("Slide subtitle")
    #   83 = footnote ("* footnote") - two-run text frame
    def fill_s2_like(slide_idx, page_num, subtitle, title, body, footnote):
        slide = prs.slides[slide_idx]
        # find shapes by their text-frame current content (since dups preserve ids)
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            t = shp.text_frame.text
            if t == "long title for slide with text":
                set_text(shp.text_frame, title)
            elif t.startswith("Lorem Ipsum is simply dummy text"):
                set_text(shp.text_frame, body)
            elif t == "02":
                set_text(shp.text_frame, page_num)
            elif t == "Slide subtitle":
                set_text(shp.text_frame, subtitle)
            elif "footnote" in t:
                set_text(shp.text_frame, footnote)
                # template footnote box is 0.89" wide → footnote text wraps
                # into a single-letter column. Widen it.
                shp.width = Inches(19.5)

    # Helper to fill a slide-3-layout-style slide (two text columns).
    # Source shape ids (slide 3):
    #   92 = page number ("03")
    #   93 = title placeholder ("long title for slide with text")
    #   94 = right-column body ("Contrary to popular belief...")
    #   113 = left-column body ("Lorem Ipsum is simply...")
    #   114 = subtitle ("Slide subtitle")
    def fill_s3_like(slide_idx, page_num, subtitle, title, left, right):
        slide = prs.slides[slide_idx]
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            t = shp.text_frame.text
            if t == "long title for slide with text":
                set_text(shp.text_frame, title)
                # template title box is 10" wide; widen so longer titles fit
                shp.width = Inches(19.5)
            elif t == "03":
                set_text(shp.text_frame, page_num)
            elif t == "Slide subtitle":
                set_text(shp.text_frame, subtitle)
            elif t.startswith("Contrary to popular belief"):
                set_text(shp.text_frame, right)
            elif t.startswith("Lorem Ipsum is simply dummy text"):
                set_text(shp.text_frame, left)

    # Slide 2 — Problem
    fill_s2_like(1, S2_NUM, S2_SUB, S2_TITLE, S2_BODY, S2_FOOT)
    # Slide 3 — Vision
    fill_s3_like(2, S3_NUM, S3_SUB, S3_TITLE, S3_LEFT, S3_RIGHT)
    # Slide 4 — What we built
    fill_s2_like(3, S4_NUM, S4_SUB, S4_TITLE, S4_BODY, S4_FOOT)
    # Slide 5 — Coordinator
    fill_s2_like(4, S5_NUM, S5_SUB, S5_TITLE, S5_BODY, S5_FOOT)
    # Slide 6 — Agents
    fill_s3_like(5, S6_NUM, S6_SUB, S6_TITLE, S6_LEFT, S6_RIGHT)
    # Slide 7 — Split-second
    fill_s2_like(6, S7_NUM, S7_SUB, S7_TITLE, S7_BODY, S7_FOOT)

    # Slide 8 — Results (slide-4 layout). Fill title + subtitle text only here;
    # chart/table data is rewritten in a separate pass below.
    slide = prs.slides[7]
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        t = shp.text_frame.text
        if t == "slide title with table and graphs":
            set_text(shp.text_frame, S8_TITLE)
        elif t == "04":
            set_text(shp.text_frame, S8_NUM)
        elif t == "Slide subtitle":
            set_text(shp.text_frame, S8_SUB)

    # Slide 9 — Today vs production
    fill_s3_like(8, S9_NUM, S9_SUB, S9_TITLE, S9_LEFT, S9_RIGHT)

    # Slide 10 — Ask
    fill_s2_like(9, S10_NUM, S10_SUB, S10_TITLE, S10_BODY, S10_FOOT)

    # ---------- slide 4: embed dashboard screenshot ----------
    # Slide 4 ("What we built") body text currently fills the right half. Move
    # the body up + shrink its height, then drop the live-sim PNG below it so
    # the slide doubles as a visual reference to what the video will show.
    slide = prs.slides[3]
    for shp in slide.shapes:
        if shp.has_text_frame and shp.text_frame.text.startswith("HORUS is not a slide deck"):
            shp.top = Inches(2.74)
            shp.height = Inches(4.0)
    # try the current filename first, then fall back to older names
    for fname in ("demo_hivemind.png", "dashboard_hivemind_roles.png",
                  "dashboard_seed1.png"):
        screenshot = ROOT / fname
        if screenshot.exists():
            break
    else:
        screenshot = None
    if screenshot is not None:
        slide.shapes.add_picture(
            str(screenshot),
            Inches(11.0), Inches(7.0),
            width=Inches(10.0),
        )

    # ---------- slide 8: rewrite table, delete charts + chart callouts ----------
    # The template's charts come with embedded Excel workbooks whose data
    # python-pptx's `replace_data` does not reliably overwrite. Rather than
    # ship a slide where the chart bars/segments contradict the table
    # numbers, we delete both charts and their doughnut callouts and let the
    # table + a sim screenshot carry the slide.
    slide = prs.slides[7]
    table_shape = None
    shapes_to_remove = []
    for shp in slide.shapes:
        if shp.has_table:
            table_shape = shp
        elif shp.has_chart:
            shapes_to_remove.append(shp)
        elif shp.has_text_frame and shp.text_frame.text.strip() in (
            "12.5%", "25%", "50%"
        ):
            shapes_to_remove.append(shp)
        elif shp.shape_type is not None and "FREEFORM" in str(shp.shape_type) and shp.width < Inches(1):
            # the small triangle/arrow shapes that were paired with the
            # doughnut callouts — only the tiny ones, not the full-width bands
            shapes_to_remove.append(shp)
    for shp in shapes_to_remove:
        el = shp._element
        el.getparent().remove(el)
    if table_shape is not None:
        rewrite_table(table_shape.table)
        # widen the table so it dominates the slide
        table_shape.left = Inches(1.33)
        table_shape.top = Inches(4.7)
        table_shape.width = Inches(11.0)
    # add a screenshot of the live hivemind sim on the right
    for fname in ("demo_hivemind.png", "dashboard_hivemind_roles.png"):
        screenshot = ROOT / fname
        if screenshot.exists():
            slide.shapes.add_picture(
                str(screenshot),
                Inches(13.0), Inches(4.7),
                width=Inches(8.5),
            )
            break

    prs.save(str(OUTPUT))
    print(f"wrote {OUTPUT}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
